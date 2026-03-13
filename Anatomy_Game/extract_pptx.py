#!/usr/bin/env python3
"""
Extract images and labels from PowerPoint slides for BabyduckGame.
Run authoring mode to define where each label should be placed on the image.
"""

from __future__ import annotations

import argparse
import json
import sys
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

try:
    from PIL import Image
    import io
except ImportError:
    Image = None
    io = None


def _collect_text_shapes(shapes, out: list[tuple[str, float, float, float, float]]):
    """Recursively collect (text, left, top, width, height) from shapes. Coords in EMUs."""
    for s in shapes:
        if s.shape_type == MSO_SHAPE_TYPE.GROUP:
            _collect_text_shapes(s.shapes, out)
            continue
        text = None
        if hasattr(s, "text") and s.text and s.text.strip():
            text = s.text.strip()
        elif hasattr(s, "has_text_frame") and s.has_text_frame and s.text_frame:
            t = s.text_frame.text.strip()
            if t:
                text = t
        if text and len(text) <= 200:
            # Use shape position; EMUs to inches: 914400 EMUs = 1 inch
            left = getattr(s, "left", None) or 0
            top = getattr(s, "top", None) or 0
            width = getattr(s, "width", None) or 0
            height = getattr(s, "height", None) or 0
            out.append((text, float(left), float(top), float(width), float(height)))


def _collect_pictures(shapes, out: list):
    """Recursively collect picture shapes."""
    for s in shapes:
        if s.shape_type == MSO_SHAPE_TYPE.GROUP:
            _collect_pictures(s.shapes, out)
            continue
        if s.shape_type == MSO_SHAPE_TYPE.PICTURE:
            out.append(s)


def extract_slide(prs: Presentation, slide_index: int) -> tuple[bytes | None, str, list[tuple[str, float, float, float, float]]]:
    """Extract image bytes (first/largest picture), extension, and label list from a slide."""
    slide = prs.slides[slide_index]
    pictures = []
    _collect_pictures(slide.shapes, pictures)
    labels = []
    _collect_text_shapes(slide.shapes, labels)

    # Prefer largest picture by area
    def area(s):
        w = getattr(s, "width", 0) or 0
        h = getattr(s, "height", 0) or 0
        return float(w) * float(h)

    pictures.sort(key=area, reverse=True)
    image_blob = None
    image_ext = "png"
    if pictures:
        img = pictures[0].image
        image_blob = img.blob
        image_ext = (img.ext or "png").lower()
        if image_ext not in ("png", "jpg", "jpeg", "gif", "bmp", "webp"):
            image_ext = "png"

    return image_blob, image_ext, labels


def extract_pptx(pptx_path: Path, output_dir: Path, slide_index: int | None = None) -> list[Path]:
    """Extract images and labels from pptx. Returns paths to created level dirs."""
    prs = Presentation(str(pptx_path))
    output_dir.mkdir(parents=True, exist_ok=True)
    level_dirs = []

    indices = [slide_index] if slide_index is not None else range(len(prs.slides))
    for i in indices:
        if i < 0 or i >= len(prs.slides):
            continue
        image_blob, image_ext, labels = extract_slide(prs, i)
        level_name = f"slide_{i}"
        level_dir = output_dir / level_name
        level_dir.mkdir(parents=True, exist_ok=True)

        image_path = level_dir / f"image.{image_ext}"
        if image_blob:
            with open(image_path, "wb") as f:
                f.write(image_blob)
        else:
            # No picture on slide; create a placeholder so authoring/game don't break
            if Image and io:
                w, h = 800, 600
                im = Image.new("RGB", (w, h), (240, 240, 240))
                buf = io.BytesIO()
                im.save(buf, "PNG")
                image_path = level_dir / "image.png"
                with open(image_path, "wb") as f:
                    f.write(buf.getvalue())
            else:
                raise RuntimeError("No image on slide and Pillow not installed for placeholder.")

        level_data = {
            "image": image_path.name,
            "labels": [{"text": t, "x": None, "y": None} for t, *_ in labels],
        }
        json_path = level_dir / "level.json"
        with open(json_path, "w", encoding="utf-8") as f:
            json.dump(level_data, f, indent=2, ensure_ascii=False)

        level_dirs.append(level_dir)
    return level_dirs


def run_authoring(level_dir: Path) -> None:
    """GUI to click where each label belongs on the image. Saves positions to level.json."""
    import tkinter as tk
    from tkinter import font as tkfont

    json_path = level_dir / "level.json"
    with open(json_path, "r", encoding="utf-8") as f:
        data = json.load(f)
    image_name = data["image"]
    image_path = level_dir / image_name
    if not image_path.exists():
        raise FileNotFoundError(f"Image not found: {image_path}")

    if not Image:
        raise RuntimeError("Pillow required for authoring. pip install Pillow")

    img = Image.open(image_path).convert("RGB")
    width, height = img.size
    labels = data["labels"]
    index = [0]

    def on_canvas_click(event):
        i = index[0]
        if i >= len(labels):
            root.quit()
            return
        # Map canvas coords to image coords (canvas may be scaled)
        cw = canvas.winfo_width()
        ch = canvas.winfo_height()
        if cw <= 1 or ch <= 1:
            return
        scale = min(cw / width, ch / height)
        off_x = (cw - width * scale) / 2
        off_y = (ch - height * scale) / 2
        ix = (event.x - off_x) / scale
        iy = (event.y - off_y) / scale
        ix = max(0, min(width, ix))
        iy = max(0, min(height, iy))
        labels[i]["x"] = round(ix, 1)
        labels[i]["y"] = round(iy, 1)
        index[0] += 1
        update_ui()

    def update_ui():
        i = index[0]
        if i < len(labels):
            status.set(f"Click where '{labels[i]['text']}' points to ({i + 1}/{len(labels)})")
        else:
            status.set("All done! Close the window.")
            with open(json_path, "w", encoding="utf-8") as f:
                json.dump(data, f, indent=2, ensure_ascii=False)

    root = tk.Tk()
    root.title("BabyduckGame - Authoring: define label positions")
    root.geometry("900x700")
    status = tk.StringVar()
    update_ui()
    tk.Label(root, textvariable=status, font=tkfont.Font(size=12)).pack(pady=8)
    canvas = tk.Canvas(root, width=800, height=600, bg="gray85")
    canvas.pack(pady=8, expand=True, fill=tk.BOTH)
    canvas.bind("<Button-1>", on_canvas_click)

    scale = min(800 / width, 600 / height)
    disp_w, disp_h = int(width * scale), int(height * scale)
    resample = getattr(Image, "Resampling", None)
    resample = (resample.LANCZOS if resample else None) or getattr(Image, "LANCZOS", Image.NEAREST)
    img_resized = img.resize((disp_w, disp_h), resample)
    try:
        from PIL import ImageTk
        photo = ImageTk.PhotoImage(img_resized)
    except Exception:
        photo = None
    if photo:
        canvas.create_image(400, 300, image=photo, tags="img")
        canvas.image = photo

    root.mainloop()


def main():
    ap = argparse.ArgumentParser(description="Extract images and labels from PowerPoint for BabyduckGame.")
    ap.add_argument("pptx", type=Path, help="Path to .pptx file")
    ap.add_argument("-o", "--output", type=Path, default=Path("levels"), help="Output directory for levels")
    ap.add_argument("-s", "--slide", type=int, default=None, help="Extract only this 0-based slide index")
    ap.add_argument("--author", type=Path, default=None, help="Run authoring for this level dir (e.g. levels/slide_0)")
    args = ap.parse_args()

    if args.author:
        run_authoring(Path(args.author))
        return

    if not args.pptx.exists():
        print(f"Not found: {args.pptx}", file=sys.stderr)
        sys.exit(1)
    dirs = extract_pptx(args.pptx, args.output, args.slide)
    print(f"Extracted {len(dirs)} level(s):")
    for d in dirs:
        print(f"  {d}")
    print("\nNext: run authoring for each level to set label positions, e.g.")
    print(f"  python extract_pptx.py {args.pptx} --author {dirs[0]}")


if __name__ == "__main__":
    main()
